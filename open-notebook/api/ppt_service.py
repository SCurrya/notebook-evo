"""
PPT 生成服务。

使用 python-pptx 生成 .pptx 文件，支持多种内置模板、Markdown 解析、
自定义标题幻灯片、图片插入（URL 或 base64）。

设计：
- PPTService 为单例（参考 agent_service.py 的 AgentManager 模式）
- 任务状态保存在内存字典中（进程内）
- 同步的 python-pptx 调用通过 asyncio.to_thread 包装为异步
- 生成的文件保存到 outputs/ppt/ 目录

模板：
- default: 简洁白底，靛紫强调色
- business: 商务深蓝，专业风格
- tech: 科技深色，青绿强调色
- education: 教育暖色，橙色强调色
- creative: 创意渐变，粉紫强调色
"""

import asyncio
import base64
import io
import uuid
from datetime import datetime, timezone
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import httpx
from loguru import logger
from pydantic import BaseModel, Field


# === 枚举 ===

class PPTTemplateType(str, Enum):
    """PPT 模板类型。"""
    DEFAULT = "default"
    BUSINESS = "business"
    TECH = "tech"
    EDUCATION = "education"
    CREATIVE = "creative"


class PPTTaskState(str, Enum):
    """PPT 任务状态。"""
    PENDING = "pending"
    PROCESSING = "processing"
    COMPLETED = "completed"
    FAILED = "failed"


# === Pydantic 模型 ===

class PPTSlide(BaseModel):
    """单张幻灯片内容。"""
    title: str = Field(..., description="幻灯片标题")
    content: List[str] = Field(
        default_factory=list,
        description="内容条目（bullet points）",
    )
    image_url: Optional[str] = Field(
        default=None, description="图片 URL（http/https）",
    )
    image_base64: Optional[str] = Field(
        default=None, description="图片 base64 编码（可含 data URI 前缀）",
    )
    layout: str = Field(
        default="title_content",
        description="布局类型: title_content / title_only / content_only / image",
    )


class PPTGenerationRequest(BaseModel):
    """PPT 生成请求。"""
    title: str = Field(..., description="演示文稿标题（用于标题幻灯片）")
    subtitle: Optional[str] = Field(default=None, description="副标题")
    author: Optional[str] = Field(default=None, description="作者")
    date: Optional[str] = Field(
        default=None,
        description="日期字符串（如不提供则使用当前日期）",
    )
    template: PPTTemplateType = Field(
        default=PPTTemplateType.DEFAULT,
        description="模板类型",
    )
    markdown_content: Optional[str] = Field(
        default=None,
        description="Markdown 内容（## 作为新幻灯片标题，- / * 作为 bullet）",
    )
    slides: Optional[List[PPTSlide]] = Field(
        default=None,
        description="显式幻灯片列表（若提供则与 markdown 合并）",
    )


class PPTTemplate(BaseModel):
    """模板信息（用于前端展示）。"""
    id: str
    name: str
    description: str
    # 前端用于渐变预览的颜色（hex 字符串）
    preview_colors: List[str]
    accent_color: str


class PPTTaskResponse(BaseModel):
    """生成任务提交响应。"""
    task_id: str
    status: str = "pending"
    message: str = ""


class PPTTaskStatus(BaseModel):
    """PPT 任务状态。"""
    id: str
    state: str
    progress: int = 0
    message: str = ""
    file_path: Optional[str] = None
    download_url: Optional[str] = None
    created_at: str
    updated_at: Optional[str] = None
    completed_at: Optional[str] = None
    error: Optional[str] = None
    template: Optional[str] = None
    title: Optional[str] = None
    slide_count: Optional[int] = None


class PPTTaskListResponse(BaseModel):
    """任务列表分页响应。"""
    items: List[PPTTaskStatus]
    total: int
    page: int
    page_size: int


# === 模板配置 ===
# 每个模板定义配色方案（RGBColor 元组）、字体和前端预览色

def _rgb(r: int, g: int, b: int) -> Tuple[int, int, int]:
    """返回 RGB 元组（0-255）。"""
    return (r, g, b)


TEMPLATE_CONFIGS: Dict[str, Dict[str, Any]] = {
    "default": {
        "name": "Default",
        "description": "简洁白底，靛紫强调色，适合通用场景",
        "colors": {
            "background": _rgb(0xFF, 0xFF, 0xFF),
            "title_bg": _rgb(0x63, 0x66, 0xF1),
            "title": _rgb(0x1F, 0x29, 0x37),
            "content": _rgb(0x4B, 0x55, 0x63),
            "accent": _rgb(0x63, 0x66, 0xF1),  # indigo-500
            "subtitle": _rgb(0x6B, 0x72, 0x80),
        },
        "fonts": {
            "title": "Calibri",
            "content": "Calibri",
        },
        "preview_colors": ["#6366F1", "#A5B4FC", "#FFFFFF"],
        "accent_color": "#6366F1",
    },
    "business": {
        "name": "Business",
        "description": "商务深蓝，专业稳重，适合企业汇报",
        "colors": {
            "background": _rgb(0xFF, 0xFF, 0xFF),
            "title_bg": _rgb(0x1E, 0x3A, 0x8A),
            "title": _rgb(0x1E, 0x3A, 0x8A),  # 深蓝
            "content": _rgb(0x37, 0x41, 0x51),
            "accent": _rgb(0x25, 0x63, 0xEB),  # blue-600
            "subtitle": _rgb(0x6B, 0x72, 0x80),
        },
        "fonts": {
            "title": "Arial",
            "content": "Arial",
        },
        "preview_colors": ["#1E3A8A", "#2563EB", "#DBEAFE"],
        "accent_color": "#2563EB",
    },
    "tech": {
        "name": "Tech",
        "description": "科技深色底，青绿强调色，适合技术分享",
        "colors": {
            "background": _rgb(0x0F, 0x17, 0x2A),  # 深色底
            "title_bg": _rgb(0x0F, 0x17, 0x2A),
            "title": _rgb(0x34, 0xD3, 0x99),  # emerald-400
            "content": _rgb(0xCB, 0xD5, 0xE1),  # slate-300
            "accent": _rgb(0x06, 0xB6, 0xD4),  # cyan-500
            "subtitle": _rgb(0x94, 0xA3, 0xB8),  # slate-400
        },
        "fonts": {
            "title": "Consolas",
            "content": "Segoe UI",
        },
        "preview_colors": ["#0F172A", "#06B6D4", "#34D399"],
        "accent_color": "#06B6D4",
    },
    "education": {
        "name": "Education",
        "description": "教育暖色底，橙色强调色，适合教学课件",
        "colors": {
            "background": _rgb(0xFF, 0xFB, 0xEB),  # 暖白
            "title_bg": _rgb(0xFF, 0xFB, 0xEB),
            "title": _rgb(0x9A, 0x34, 0x12),  # orange-800
            "content": _rgb(0x44, 0x40, 0x3C),
            "accent": _rgb(0xF5, 0x9E, 0x0B),  # amber-500
            "subtitle": _rgb(0x78, 0x71, 0x6C),
        },
        "fonts": {
            "title": "Georgia",
            "content": "Verdana",
        },
        "preview_colors": ["#FFF8EB", "#F59E0B", "#9A3412"],
        "accent_color": "#F59E0B",
    },
    "creative": {
        "name": "Creative",
        "description": "创意粉紫，活泼灵动，适合创意提案",
        "colors": {
            "background": _rgb(0xFD, 0xF4, 0xFF),  # 淡紫底
            "title_bg": _rgb(0xFD, 0xF4, 0xFF),
            "title": _rgb(0x86, 0x16, 0x98),  # purple-800
            "content": _rgb(0x44, 0x40, 0x3C),
            "accent": _rgb(0xEC, 0x48, 0x99),  # pink-500
            "subtitle": _rgb(0xA2, 0x1C, 0xAF),
        },
        "fonts": {
            "title": "Trebuchet MS",
            "content": "Trebuchet MS",
        },
        "preview_colors": ["#FDF4FF", "#EC4899", "#861698"],
        "accent_color": "#EC4899",
    },
}


# === PPTService ===

class PPTService:
    """PPT 生成服务（单例）。

    管理生成任务的生命周期：提交 → 处理 → 完成/失败。
    所有状态保存在内存中（进程内），参考 agent_service.py 模式。
    """

    _instance: Optional["PPTService"] = None

    def __init__(self) -> None:
        self._tasks: Dict[str, PPTTaskStatus] = {}
        self._task_requests: Dict[str, PPTGenerationRequest] = {}
        # 输出目录：项目根 / outputs / ppt
        self.output_dir: Path = Path(__file__).parent.parent / "outputs" / "ppt"
        self.output_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"PPTService initialized, output_dir={self.output_dir}")

    @classmethod
    def get_instance(cls) -> "PPTService":
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    # --- 模板 ---

    @staticmethod
    def list_templates() -> List[PPTTemplate]:
        """列出所有可用模板。"""
        templates: List[PPTTemplate] = []
        for tid, cfg in TEMPLATE_CONFIGS.items():
            templates.append(
                PPTTemplate(
                    id=tid,
                    name=cfg["name"],
                    description=cfg["description"],
                    preview_colors=cfg["preview_colors"],
                    accent_color=cfg["accent_color"],
                )
            )
        return templates

    # --- 任务管理 ---

    async def create_task(
        self, request: PPTGenerationRequest
    ) -> PPTTaskResponse:
        """提交 PPT 生成任务。

        立即返回 task_id，生成在后台异步执行。
        """
        task_id = str(uuid.uuid4())
        now = datetime.now(timezone.utc).isoformat()

        task = PPTTaskStatus(
            id=task_id,
            state=PPTTaskState.PENDING.value,
            progress=0,
            message="Task queued",
            created_at=now,
            updated_at=now,
            template=request.template.value,
            title=request.title,
        )
        self._tasks[task_id] = task
        self._task_requests[task_id] = request

        logger.info(
            f"PPT task created: id={task_id} title='{request.title}' "
            f"template={request.template}"
        )

        # 后台异步执行生成
        asyncio.create_task(self._run_generation(task_id, request))

        return PPTTaskResponse(
            task_id=task_id,
            status="pending",
            message="PPT generation task submitted",
        )

    def get_task(self, task_id: str) -> Optional[PPTTaskStatus]:
        """获取单个任务状态。"""
        return self._tasks.get(task_id)

    def list_tasks(
        self, page: int = 1, page_size: int = 20
    ) -> PPTTaskListResponse:
        """分页列出任务（按创建时间降序）。"""
        all_tasks = list(self._tasks.values())
        all_tasks.sort(key=lambda t: t.created_at, reverse=True)
        total = len(all_tasks)
        start = (page - 1) * page_size
        end = start + page_size
        items = all_tasks[start:end]
        return PPTTaskListResponse(
            items=items,
            total=total,
            page=page,
            page_size=page_size,
        )

    def delete_task(self, task_id: str) -> bool:
        """删除任务及其生成的文件。"""
        task = self._tasks.get(task_id)
        if not task:
            return False

        # 删除文件
        if task.file_path:
            file_path = Path(task.file_path)
            if file_path.exists():
                try:
                    file_path.unlink()
                    logger.info(f"Deleted PPT file: {file_path}")
                except OSError as e:
                    logger.warning(f"Failed to delete PPT file {file_path}: {e}")

        del self._tasks[task_id]
        self._task_requests.pop(task_id, None)
        logger.info(f"PPT task deleted: {task_id}")
        return True

    def get_file_path(self, task_id: str) -> Optional[Path]:
        """获取任务生成文件的路径（用于下载）。"""
        task = self._tasks.get(task_id)
        if not task or not task.file_path:
            return None
        path = Path(task.file_path)
        return path if path.exists() else None

    # --- 异步生成 ---

    async def _run_generation(
        self, task_id: str, request: PPTGenerationRequest
    ) -> None:
        """后台执行 PPT 生成（包装同步 python-pptx 调用）。"""
        task = self._tasks.get(task_id)
        if not task:
            return

        task.state = PPTTaskState.PROCESSING.value
        task.progress = 10
        task.message = "Generating presentation..."
        task.updated_at = datetime.now(timezone.utc).isoformat()

        try:
            # 通过 asyncio.to_thread 包装同步生成逻辑
            file_path = await asyncio.to_thread(
                self._generate_pptx_sync, request, task_id
            )

            task.file_path = file_path
            task.download_url = f"/api/ppt/tasks/{task_id}/download"
            task.progress = 100
            task.state = PPTTaskState.COMPLETED.value
            task.message = "Presentation generated successfully"
            task.completed_at = datetime.now(timezone.utc).isoformat()
            task.updated_at = task.completed_at

            # 计算幻灯片数量
            slides = self._collect_slides(request)
            task.slide_count = len(slides) + 1  # +1 for title slide

            logger.info(f"PPT task completed: id={task_id} file={file_path}")
        except Exception as e:
            task.state = PPTTaskState.FAILED.value
            task.error = str(e)
            task.message = "Generation failed"
            task.updated_at = datetime.now(timezone.utc).isoformat()
            logger.error(f"PPT task failed: id={task_id} - {e}")
            logger.exception(e)

    # --- 同步生成逻辑（python-pptx） ---

    def _generate_pptx_sync(
        self, request: PPTGenerationRequest, task_id: str
    ) -> str:
        """同步生成 PPTX 文件（在 to_thread 中调用）。"""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        template_cfg = TEMPLATE_CONFIGS.get(
            request.template.value, TEMPLATE_CONFIGS["default"]
        )

        prs = Presentation()
        # 16:9 宽屏
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 标题幻灯片
        self._add_title_slide(prs, request, template_cfg)

        # 内容幻灯片
        slides = self._collect_slides(request)
        for slide_data in slides:
            self._add_content_slide(prs, slide_data, template_cfg)

        # 保存
        output_path = self.output_dir / f"{task_id}.pptx"
        prs.save(str(output_path))
        logger.info(f"PPTX saved: {output_path}")
        return str(output_path)

    def _collect_slides(
        self, request: PPTGenerationRequest
    ) -> List[PPTSlide]:
        """收集所有内容幻灯片（合并显式 slides 和 markdown 解析结果）。"""
        slides: List[PPTSlide] = []
        if request.slides:
            slides.extend(request.slides)
        if request.markdown_content:
            slides.extend(self._parse_markdown(request.markdown_content))
        return slides

    def _add_title_slide(
        self,
        prs: Any,
        request: PPTGenerationRequest,
        cfg: Dict[str, Any],
    ) -> None:
        """添加标题幻灯片。"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        self._set_background(slide, cfg["colors"]["background"])

        # 标题
        title_left = Inches(1)
        title_top = Inches(2.2)
        title_width = Inches(11.333)
        title_height = Inches(1.5)
        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p = title_tf.paragraphs[0]
        p.text = request.title
        p.alignment = PP_ALIGN.CENTER
        self._style_font(
            p, cfg["fonts"]["title"], Pt(44), True, cfg["colors"]["title"]
        )

        # 副标题
        if request.subtitle:
            sub_top = Inches(3.8)
            sub_box = slide.shapes.add_textbox(
                Inches(1), sub_top, Inches(11.333), Inches(0.8)
            )
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sp = sub_tf.paragraphs[0]
            sp.text = request.subtitle
            sp.alignment = PP_ALIGN.CENTER
            self._style_font(
                sp, cfg["fonts"]["content"], Pt(24), False,
                cfg["colors"]["subtitle"]
            )

        # 作者 + 日期
        meta_parts: List[str] = []
        if request.author:
            meta_parts.append(request.author)
        date_str = request.date or datetime.now().strftime("%Y-%m-%d")
        meta_parts.append(date_str)

        meta_top = Inches(5.5)
        meta_box = slide.shapes.add_textbox(
            Inches(1), meta_top, Inches(11.333), Inches(0.6)
        )
        meta_tf = meta_box.text_frame
        mp = meta_tf.paragraphs[0]
        mp.text = "  ·  ".join(meta_parts)
        mp.alignment = PP_ALIGN.CENTER
        self._style_font(
            mp, cfg["fonts"]["content"], Pt(16), False, cfg["colors"]["subtitle"]
        )

        # 装饰性强调线
        self._add_accent_line(slide, cfg)

    def _add_content_slide(
        self,
        prs: Any,
        slide_data: PPTSlide,
        cfg: Dict[str, Any],
    ) -> None:
        """添加内容幻灯片。"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        self._set_background(slide, cfg["colors"]["background"])

        layout = slide_data.layout

        # 标题区域（除 content_only 外都显示标题）
        if layout != "content_only":
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0)
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            tp = title_tf.paragraphs[0]
            tp.text = slide_data.title
            tp.alignment = PP_ALIGN.LEFT
            self._style_font(
                tp, cfg["fonts"]["title"], Pt(32), True, cfg["colors"]["title"]
            )

        # 图片处理
        image_bytes: Optional[bytes] = None
        if slide_data.image_url:
            image_bytes = self._download_image(slide_data.image_url)
        elif slide_data.image_base64:
            image_bytes = self._decode_base64_image(slide_data.image_base64)

        # 纯图片布局
        if layout == "image" or (image_bytes and not slide_data.content):
            if image_bytes:
                self._add_image_centered(slide, image_bytes)
            return

        # 内容区域
        if slide_data.content:
            content_top = Inches(1.8) if layout != "content_only" else Inches(0.8)
            content_height = Inches(5.2)

            # 如果有图片，内容区域缩窄
            if image_bytes:
                content_width = Inches(6.5)
                content_left = Inches(0.8)
                self._add_image_right(slide, image_bytes)
            else:
                content_width = Inches(11.733)
                content_left = Inches(0.8)

            content_box = slide.shapes.add_textbox(
                content_left, content_top, content_width, content_height
            )
            content_tf = content_box.text_frame
            content_tf.word_wrap = True

            for idx, item in enumerate(slide_data.content):
                if idx == 0:
                    p = content_tf.paragraphs[0]
                else:
                    p = content_tf.add_paragraph()
                p.text = f"• {item}"
                p.space_after = Pt(8)
                self._style_font(
                    p, cfg["fonts"]["content"], Pt(20), False,
                    cfg["colors"]["content"]
                )
        elif image_bytes and layout != "title_only":
            # 有图片但无文字内容
            self._add_image_centered(slide, image_bytes)

    # --- 辅助方法 ---

    @staticmethod
    def _set_background(slide: Any, rgb: Tuple[int, int, int]) -> None:
        """设置幻灯片背景色。"""
        from pptx.dml.color import RGBColor

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)

    @staticmethod
    def _style_font(
        paragraph: Any,
        font_name: str,
        size: Any,
        bold: bool,
        rgb: Tuple[int, int, int],
    ) -> None:
        """设置段落字体样式。"""
        from pptx.dml.color import RGBColor

        font = paragraph.font
        font.name = font_name
        font.size = size
        font.bold = bold
        font.color.rgb = RGBColor(*rgb)

    @staticmethod
    def _add_accent_line(slide: Any, cfg: Dict[str, Any]) -> None:
        """在标题幻灯片底部添加装饰性强调线。"""
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        rgb = cfg["colors"]["accent"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.667),  # 居中
            Inches(4.8),
            Inches(2.0),
            Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*rgb)
        shape.line.fill.background()  # 无边框

    @staticmethod
    def _add_image_centered(slide: Any, image_bytes: bytes) -> None:
        """将图片居中添加到幻灯片。"""
        from pptx.util import Inches

        image_stream = io.BytesIO(image_bytes)
        # 居中放置，限制最大尺寸
        left = Inches(3.0)
        top = Inches(1.8)
        width = Inches(7.333)
        height = Inches(4.5)
        try:
            slide.shapes.add_picture(
                image_stream, left, top, width=width, height=height
            )
        except Exception as e:
            logger.warning(f"Failed to add centered image: {e}")

    @staticmethod
    def _add_image_right(slide: Any, image_bytes: bytes) -> None:
        """将图片添加到幻灯片右侧。"""
        from pptx.util import Inches

        image_stream = io.BytesIO(image_bytes)
        left = Inches(7.8)
        top = Inches(1.8)
        width = Inches(4.733)
        height = Inches(4.5)
        try:
            slide.shapes.add_picture(
                image_stream, left, top, width=width, height=height
            )
        except Exception as e:
            logger.warning(f"Failed to add right image: {e}")

    @staticmethod
    def _download_image(url: str) -> bytes:
        """从 URL 下载图片（同步，在 to_thread 中调用）。"""
        logger.info(f"Downloading image: {url}")
        with httpx.Client(timeout=30.0, follow_redirects=True) as client:
            response = client.get(url)
            response.raise_for_status()
            return response.content

    @staticmethod
    def _decode_base64_image(b64: str) -> bytes:
        """解码 base64 图片（可含 data URI 前缀）。"""
        # 移除 data URI 前缀（如 data:image/png;base64,...）
        if "," in b64 and b64.startswith("data:"):
            b64 = b64.split(",", 1)[1]
        return base64.b64decode(b64)

    @staticmethod
    def _parse_markdown(markdown: str) -> List[PPTSlide]:
        """将 Markdown 内容解析为幻灯片列表。

        解析规则：
        - `## 标题` 或 `# 标题`：开始新幻灯片
        - `- item` / `* item` / `+ item`：bullet point
        - `1. item`：编号列表（转为 bullet）
        - 其他文本：作为普通内容行
        - `![alt](url)`：图片 URL
        """
        slides: List[PPTSlide] = []
        current_title: Optional[str] = None
        current_content: List[str] = []
        current_image_url: Optional[str] = None

        def _flush() -> None:
            nonlocal current_title, current_content, current_image_url
            if current_title is not None:
                slide = PPTSlide(
                    title=current_title,
                    content=current_content,
                    image_url=current_image_url,
                )
                slides.append(slide)
            current_title = None
            current_content = []
            current_image_url = None

        for raw_line in markdown.splitlines():
            line = raw_line.strip()
            if not line:
                continue

            # 标题（## 或 #）→ 新幻灯片
            if line.startswith("## "):
                _flush()
                current_title = line[3:].strip()
            elif line.startswith("# "):
                _flush()
                current_title = line[2:].strip()
            elif line.startswith("### "):
                # 三级标题作为内容子标题
                current_content.append(line[4:].strip())
            elif line.startswith(("- ", "* ", "+ ")):
                current_content.append(line[2:].strip())
            elif line.startswith("![") and "](" in line:
                # Markdown 图片语法 ![alt](url)
                url_start = line.find("](") + 2
                url_end = line.rfind(")")
                if url_start > 1 and url_end > url_start:
                    current_image_url = line[url_start:url_end]
            else:
                # 编号列表或其他文本
                # 检测 "1. " 等编号列表
                dot_idx = line.find(". ")
                if (
                    dot_idx > 0
                    and line[:dot_idx].isdigit()
                ):
                    current_content.append(line[dot_idx + 2:].strip())
                else:
                    current_content.append(line)

        _flush()
        return slides
